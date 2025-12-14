from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(output_file="AI_DevOps_Analyst_Presentation.pptx"):
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI DevOps Incident Analyst"
    subtitle.text = "Automated Incident Remediation with Multi-Agent Architecture\nPresented by: AI Assistant"

    # Slide 2: Agenda
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Agenda"
    tf = body_shape.text_frame
    tf.text = "Problem Statement"
    p = tf.add_paragraph()
    p.text = "Solution Overview (AI DevOps Analyst)"
    p = tf.add_paragraph()
    p.text = "Architecture & Design"
    p = tf.add_paragraph()
    p.text = "Key Benefits"
    p = tf.add_paragraph()
    p.text = "Demo/Walkthrough"

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Challenge"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "High volume of operational logs and alerts."
    p = tf.add_paragraph()
    p.text = "Slow manual analysis and correlation."
    p = tf.add_paragraph()
    p.text = "Inconsistent remediation steps."
    p = tf.add_paragraph()
    p.text = "Delayed communication to stakeholders."

    # Slide 4: Solution Overview
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "The Solution"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "AI-Driven Incident Analysis Application."
    p = tf.add_paragraph()
    p.text = "Multi-Agent System for specialized tasks (Analysis, Research, Fix)."
    p = tf.add_paragraph()
    p.text = "Integration with existing tools (JIRA, Slack)."
    p = tf.add_paragraph()
    p.text = "RAG-powered Context from Internal Cookbooks."

    # Slide 5: Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Architecture - Multi-Agent Flow"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Orchestrator: LangGraph State Machine"
    p = tf.add_paragraph()
    p.level = 1
    p.text = "1. Log Reader Agent: Parses raw logs."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "2. Cookbook Agent: Retrieves vector-embedded docs."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "3. Remediation Agent: Synthesizes fixes."
    p = tf.add_paragraph()
    p.level = 1
    p.text = "4. Action Agents: JIRA Ticket Creation & Slack Notification."

    # Slide 6: Benefits
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key Benefits"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Speed: Instant analysis of complex logs."
    p = tf.add_paragraph()
    p.text = "Accuracy: Context-aware fixes using internal knowledge bases."
    p = tf.add_paragraph()
    p.text = "Scalability: Handle thousands of incidents automatically."
    p = tf.add_paragraph()
    p.text = "Traceability: All actions logged and tracked in JIRA."

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()

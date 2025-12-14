"""
Multi-Agent AI Deep Researcher - ADVANCED VERSION
Features: RAG, PDF Upload, Citations, Streaming, Agent Visualization

Install: pip install streamlit langchain langchain-groq langgraph tavily-python python-dotenv chromadb pypdf sentence-transformers graphviz pydot
Run: streamlit run app.py
"""

import streamlit as st
import os
from dotenv import load_dotenv
from typing import TypedDict, Annotated, List, Dict, Any
import operator
import json
import re
import time
from datetime import datetime
import base64
from io import BytesIO

# LangChain & LangGraph
from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage
from langchain_core.documents import Document
from langgraph.graph import StateGraph, END
from tavily import TavilyClient

# RAG Components
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
# from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_text_splitters import RecursiveCharacterTextSplitter  # ✅ CORRECT
# Replace these imports for 3.11 version
from langchain_community.document_loaders import PyPDFLoader  # ← NEW
from langchain_community.vectorstores import Chroma           # ← NEW  
from langchain_community.embeddings import HuggingFaceEmbeddings  # ← NEW


# PDF Processing
from pypdf import PdfReader

load_dotenv()


# Page config
st.set_page_config(
    page_title="MAAD Researcher",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ===== SESSION STATE - MUST COME FIRST =====
if 'results' not in st.session_state: st.session_state.results = None
if 'agent_logs' not in st.session_state: st.session_state.agent_logs = []
if 'current_query' not in st.session_state: st.session_state.current_query = ""
if 'citations' not in st.session_state: st.session_state.citations = []
if 'vectorstore' not in st.session_state: st.session_state.vectorstore = None
if 'uploaded_docs' not in st.session_state: st.session_state.uploaded_docs = []
if 'agent_timeline' not in st.session_state: st.session_state.agent_timeline = []

# SEARCH SCOPE (CRITICAL - ADD THESE)
if 'search_scope' not in st.session_state:
    st.session_state.search_scope = "both"
if 'web_domains' not in st.session_state:
    st.session_state.web_domains = []
if 'search_history' not in st.session_state:
    st.session_state.search_history = []


# Initialize session state
if 'results' not in st.session_state:
    st.session_state.results = None
if 'agent_logs' not in st.session_state:
    st.session_state.agent_logs = []
if 'current_query' not in st.session_state:
    st.session_state.current_query = ""
if 'citations' not in st.session_state:
    st.session_state.citations = []
if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = None
if 'uploaded_docs' not in st.session_state:
    st.session_state.uploaded_docs = []
if 'agent_timeline' not in st.session_state:
    st.session_state.agent_timeline = []
# SEARCH HISTORY
if 'search_history' not in st.session_state:
    st.session_state.search_history = []  # List of past searches
if 'selected_history_id' not in st.session_state:
    st.session_state.selected_history_id = None

if 'plagiarism_active' not in st.session_state:
    st.session_state.plagiarism_active = False
if 'plagiarism_result' not in st.session_state:
    st.session_state.plagiarism_result = None



import requests

# ===== WINSTON AI PLAGIARISM CHECKER =====
def check_plagiarism(text: str, api_key: str) -> dict:
    """Check plagiarism using Winston AI API"""
    url = "https://api.gowinston.ai/v2/plagiarism"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "text": text,
        "type": "ORIGINALITY"
    }
    
    try:
        response = requests.post(url, json=payload, headers=headers)
        if response.status_code == 200:
            return response.json()
        else:
            return {"error": f"API Error: {response.status_code}"}
    except Exception as e:
        return {"error": str(e)}

@st.cache_data(ttl=3600)  # Cache 1 hour
def analyze_plagiarism_result(result: dict, word_count: int) -> str:
    """Format plagiarism result"""
    if 'error' in result:
        return f"❌ Error: {result['error']}"
    
    ai_score = result.get('ai_score', 0) * 100
    human_score = 100 - ai_score
    
    if word_count > 500:
        return """
## 🚫 Premium Required

**Word Count**: {word_count:,} words

**Upgrade to Winston AI Premium** for documents > 500 words.

[Get Premium →](https://gowinston.ai/pricing)
        """.format(word_count=word_count)
    
    return f"""
## ✅ Plagiarism Analysis Complete

**Word Count**: {word_count:,} words

**🤖 Plagiarism**: {ai_score:.1f}%
**👤 Original Content**: {human_score:.1f}%

**Verdict**: {'🟢 Original' if ai_score < 20 else '🟡 Mixed' if ai_score < 50 else '🔴 Likely AI'}

    """


# ===== FIXED EXPORT FUNCTIONS (Replace lines 60-140) =====
import io
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
from pptx.util import Inches

def generate_short_summary(full_report: str, query: str) -> str:
    """Static short summary - no LLM needed"""
    return f"""🚀 EXECUTIVE SUMMARY: {query}

{full_report[:450]}... 
Key Finding: Critical insights from multi-agent AI analysis.
Recommendation: Implement primary recommendation immediately."""

def generate_medium_summary(full_report: str, query: str) -> str:
    """Static medium summary - no LLM needed"""
    findings = results.get('analysis', {}).get('keyfindings', ['Analysis complete'])[0] if results else "Key insights generated"
    return f"""EXECUTIVE SUMMARY: {query}

**1. KEY FINDINGS**
{findings[:300]}...

**2. CRITICAL INSIGHTS** 
Multi-agent analysis reveals high-confidence trends from {len(results.get('searchresults', []))} sources.

**3. RECOMMENDATIONS**  
Immediate action: {results.get('insights', {}).get('recommendations', ['Review full report'])[0][:200]}..."""

def export_to_pdf(content: str, title: str):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [
        Paragraph(title, styles['Title']),
        Spacer(1, 12),
        Paragraph(content, styles['Normal'])
    ]
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

def export_to_pptx(content: str, title: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = content[:250] + "..."
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# Initialize LLM and tools + Winston AI
@st.cache_resource
def init_tools():
    """Initialize Groq LLM, Tavily, and Winston AI"""
    # Groq LLM
    llm = ChatGroq(
        model="llama-3.3-70b-versatile",
        groq_api_key=os.getenv("GROQ_API_KEY"),
        temperature=0.7,
        max_tokens=2000
    )
    
    # Tavily Search
    tavily = TavilyClient(api_key=os.getenv("TAVILY_API_KEY"))
    
    # Winston AI Key (validate)
    winston_api_key = os.getenv("WINSTON_API_KEY")
    
    return llm, tavily, winston_api_key

try:
    llm, tavily_client, winston_api_key = init_tools()
    
    # Check all APIs
    api_checks = []
    if os.getenv("GROQ_API_KEY"):
        api_checks.append("✅ Groq")
    else:
        api_checks.append("❌ Groq")
    
    if os.getenv("TAVILY_API_KEY"):
        api_checks.append("✅ Tavily")
    else:
        api_checks.append("❌ Tavily")
    
    if os.getenv("WINSTON_API_KEY"):
        api_checks.append("✅ Winston AI")
    else:
        api_checks.append("⚠ Winston AI")
    
    api_status = " | ".join(api_checks)
    
except Exception as e:
    st.error(f"❌ API Connection Error: {str(e)}")
    st.stop()


# Initialize embeddings for RAG
@st.cache_resource
def init_embeddings():
    """Initialize embedding model for RAG"""
    return HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        model_kwargs={'device': 'cpu'}
    )

embeddings = init_embeddings()

# ===== VECTOR DB METRICS =====
def get_vector_db_stats():
    """Get ChromaDB collection statistics"""
    if st.session_state.vectorstore is None:
        return {
            'total_chunks': 0,
            'total_size_mb': 0,
            'dimensions': 384,  # MiniLM default
            'avg_chunk_length': 0
        }
    
    try:
        collection = st.session_state.vectorstore._collection
        total_chunks = collection.count()
        total_size_bytes = collection.count() * 384 * 4  # approx 384-dim float32
        total_size_mb = round(total_size_bytes / (1024*1024), 1)
        
        # Sample chunk lengths
        sample_docs = st.session_state.vectorstore.similarity_search("sample", k=min(10, total_chunks))
        avg_chunk_length = round(sum(len(doc.page_content) for doc in sample_docs) / max(len(sample_docs), 1))
        
        return {
            'total_chunks': total_chunks,
            'total_size_mb': total_size_mb,
            'dimensions': 384,
            'avg_chunk_length': avg_chunk_length
        }
    except:
        return {'total_chunks': 0, 'total_size_mb': 0, 'dimensions': 384, 'avg_chunk_length': 0}

def get_token_estimate(text: str) -> int:
    """Rough token count estimate (4 chars per token)"""
    return round(len(text) / 4)

# Global stats
db_stats = get_vector_db_stats()



# ============= CITATION TRACKER =============

class CitationTracker:
    """Tracks citations across agents"""
    
    def __init__(self):
        self.citations = []
    
    def add_citation(self, agent_name: str, source_title: str, source_url: str, claim: str):
        """Add a citation"""
        self.citations.append({
            'agent': agent_name,
            'source_title': source_title,
            'source_url': source_url,
            'claim': claim,
            'timestamp': datetime.now().isoformat()
        })
    
    def get_citations_by_agent(self, agent_name: str):
        """Get all citations from specific agent"""
        return [c for c in self.citations if c['agent'] == agent_name]
    
    def get_all_citations(self):
        """Get all citations"""
        return self.citations


# ============= PDF PROCESSOR =============

def process_pdf(uploaded_file):
    """Extract text from PDF and add to vector store"""
    try:
        pdf_reader = PdfReader(uploaded_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        
        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        chunks = text_splitter.split_text(text)
        
        # Create documents
        documents = [
            Document(
                page_content=chunk,
                metadata={
                    'source': uploaded_file.name,
                    'page': i // 3,  # Approximate page number
                    'type': 'uploaded_pdf'
                }
            )
            for i, chunk in enumerate(chunks)
        ]
        
        return documents, text[:500]  # Return docs and preview
        
    except Exception as e:
        st.error(f"Error processing PDF: {str(e)}")
        return [], ""


# ============= RAG HELPER =============

def add_to_vectorstore(documents: List[Document]):
    """Add documents to vector store"""
    if st.session_state.vectorstore is None:
        # Create new vector store
        st.session_state.vectorstore = Chroma.from_documents(
            documents=documents,
            embedding=embeddings,
            collection_name="research_docs"
        )
    else:
        # Add to existing vector store
        st.session_state.vectorstore.add_documents(documents)


def query_vectorstore(query: str, k: int = 3):
    """Query vector store for relevant documents"""
    if st.session_state.vectorstore is None:
        return []
    
    try:
        results = st.session_state.vectorstore.similarity_search(query, k=k)
        return results
    except:
        return []


# ============= AGENT STATE =============

class AgentState(TypedDict):
    """State shared between all agents"""
    query: str
    research_plan: Dict[str, Any]
    search_results: List[Dict[str, Any]]
    rag_results: List[Dict[str, Any]]
    analysis: Dict[str, Any]
    insights: Dict[str, Any]
    final_report: str
    messages: Annotated[List[str], operator.add]
    citations: List[Dict[str, Any]]
    timeline: Annotated[List[Dict[str, Any]], operator.add]


# ============= AGENT DEFINITIONS =============

class ResearchPlannerAgent:
    """Agent 1: Plans the research strategy"""
    
    def __init__(self, llm):
        self.llm = llm
        self.name = "Research Planner"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Create research plan"""
        start_time = time.time()
        
        # Check if RAG has relevant content
        rag_context = ""
        if st.session_state.vectorstore:
            rag_docs = query_vectorstore(state['query'], k=2)
            if rag_docs:
                rag_context = f"\n\nUploaded Documents Context:\n{rag_docs[0].page_content[:300]}"
        
        prompt = f"""You are a Research Planning Agent. Create a comprehensive research strategy.

Query: {state['query']}{rag_context}

Create a JSON response with:
1. "topics": List of 3-5 key research topics
2. "search_queries": List of 5-7 specific search queries
3. "source_types": Types of sources needed
4. "research_depth": Depth level

Return ONLY valid JSON.

Example:
{{
    "topics": ["Topic 1", "Topic 2"],
    "search_queries": ["Query 1", "Query 2"],
    "source_types": ["academic", "news"],
    "research_depth": "comprehensive"
}}"""

        try:
            response = self.llm.invoke([HumanMessage(content=prompt)])
            response_text = response.content
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            
            if json_match:
                research_plan = json.loads(json_match.group())
            else:
                research_plan = {
                    "topics": [state['query']],
                    "search_queries": [state['query']],
                    "source_types": ["web"],
                    "research_depth": "moderate"
                }
            
            state['research_plan'] = research_plan
            state['messages'].append(f"✓ Research plan created with {len(research_plan['search_queries'])} queries")
            
            # Add to timeline
            state['timeline'].append({
                'agent': self.name,
                'action': 'Planning complete',
                'duration': round(time.time() - start_time, 2),
                'timestamp': datetime.now().isoformat()
            })
            
        except Exception as e:
            state['messages'].append(f"⚠ Planning error: {str(e)}")
            state['research_plan'] = {
                "topics": [state['query']],
                "search_queries": [state['query']],
                "source_types": ["web"],
                "research_depth": "basic"
            }
        
        return state


class RAGSearchAgent:
    """Agent 1.5: Search uploaded documents"""
    
    def __init__(self):
        self.name = "RAG Search"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Search vector store for relevant content"""
        start_time = time.time()
        
        if st.session_state.vectorstore is None:
            state['rag_results'] = []
            state['messages'].append("ℹ No uploaded documents to search")
            return state
        
        try:
            # Search uploaded documents
            rag_docs = query_vectorstore(state['query'], k=5)
            
            rag_results = []
            for doc in rag_docs:
                rag_results.append({
                    'content': doc.page_content,
                    'source': doc.metadata.get('source', 'Unknown'),
                    'page': doc.metadata.get('page', 0),
                    'type': 'uploaded_document'
                })
            
            state['rag_results'] = rag_results
            state['messages'].append(f"✓ Found {len(rag_results)} relevant sections in uploaded documents")
            
            # Add to timeline
            state['timeline'].append({
                'agent': self.name,
                'action': f'Retrieved {len(rag_results)} document chunks',
                'duration': round(time.time() - start_time, 2),
                'timestamp': datetime.now().isoformat()
            })
            
        except Exception as e:
            state['messages'].append(f"⚠ RAG search error: {str(e)}")
            state['rag_results'] = []
        
        return state


class WebSearchAgent:
    """Agent 2: Searches the web using Tavily"""
    
    def __init__(self, tavily_client):
        self.tavily = tavily_client
        self.name = "Web Search"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Execute web searches"""
        start_time = time.time()
        
        all_results = []
        search_queries = state['research_plan'].get('search_queries', [state['query']])
        
        # Progress indicators
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        citation_tracker = CitationTracker()
        
        for i, query in enumerate(search_queries[:5]):
            status_text.text(f"🔍 Searching: {query[:60]}...")
            
            try:
                results = self.tavily.search(
                    query=query,
                    max_results=3,
                    search_depth="advanced"
                )
                
                for result in results.get('results', []):
                    result_data = {
                        'query': query,
                        'title': result.get('title', 'Untitled'),
                        'url': result.get('url', ''),
                        'content': result.get('content', '')[:1000],
                        'score': result.get('score', 0)
                    }
                    all_results.append(result_data)
                    
                    # Track citation
                    citation_tracker.add_citation(
                        agent_name=self.name,
                        source_title=result_data['title'],
                        source_url=result_data['url'],
                        claim=f"Retrieved for query: {query}"
                    )
                    
            except Exception as e:
                state['messages'].append(f"⚠ Search failed: {str(e)}")
            
            progress_bar.progress((i + 1) / min(len(search_queries), 5))
        
        progress_bar.empty()
        status_text.empty()
        
        state['search_results'] = all_results
        state['citations'].extend(citation_tracker.get_all_citations())
        state['messages'].append(f"✓ Found {len(all_results)} sources from Tavily")
        
        # Add to timeline
        state['timeline'].append({
            'agent': self.name,
            'action': f'Retrieved {len(all_results)} web sources',
            'duration': round(time.time() - start_time, 2),
            'timestamp': datetime.now().isoformat()
        })
        
        return state


class AnalysisAgent:
    """Agent 3: Analyzes search results critically"""
    
    def __init__(self, llm):
        self.llm = llm
        self.name = "Critical Analysis"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Analyze findings"""
        start_time = time.time()
        
        # Combine web and RAG results
        web_context = "\n\n".join([
            f"Source {i+1}: {r['title']}\n{r['content'][:400]}"
            for i, r in enumerate(state['search_results'][:10])
        ])
        
        rag_context = ""
        if state.get('rag_results'):
            rag_context = "\n\nUploaded Documents:\n" + "\n\n".join([
                f"Doc {i+1}: {r['content'][:300]}"
                for i, r in enumerate(state['rag_results'][:3])
            ])
        
        prompt = f"""You are a Critical Analysis Agent.

Query: {state['query']}

Web Sources:
{web_context}
{rag_context}

Provide JSON analysis with:
1. "key_findings": 5-7 findings
2. "contradictions": conflicts found
3. "data_quality": reliability assessment
4. "confidence_level": "low", "medium", or "high"
5. "gaps": missing information

Return ONLY valid JSON."""

        try:
            response = self.llm.invoke([HumanMessage(content=prompt)])
            json_match = re.search(r'\{.*\}', response.content, re.DOTALL)
            
            if json_match:
                analysis = json.loads(json_match.group())
            else:
                analysis = {
                    "key_findings": ["Analysis completed"],
                    "contradictions": "None identified",
                    "data_quality": "Moderate",
                    "confidence_level": "medium",
                    "gaps": "Limited data"
                }
            
            state['analysis'] = analysis
            state['messages'].append(f"✓ Analysis complete: {len(analysis.get('key_findings', []))} findings")
            
            # Add to timeline
            state['timeline'].append({
                'agent': self.name,
                'action': 'Analysis completed',
                'duration': round(time.time() - start_time, 2),
                'timestamp': datetime.now().isoformat()
            })
            
        except Exception as e:
            state['messages'].append(f"⚠ Analysis error: {str(e)}")
            state['analysis'] = {
                "key_findings": ["Error"],
                "contradictions": "Unknown",
                "data_quality": "Unknown",
                "confidence_level": "low",
                "gaps": "Analysis incomplete"
            }
        
        return state


class InsightGenerationAgent:
    """Agent 4: Generates insights"""
    
    def __init__(self, llm):
        self.llm = llm
        self.name = "Insight Generator"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Generate insights"""
        start_time = time.time()
        
        prompt = f"""You are an Insight Generation Agent.

Query: {state['query']}
Findings: {state['analysis'].get('key_findings', [])}

Generate JSON with:
1. "main_insights": 3-5 takeaways
2. "trends": 2-4 patterns
3. "implications": 2-3 implications
4. "recommendations": 3-5 recommendations
5. "future_research": 2-3 areas

Return ONLY valid JSON."""

        try:
            response = self.llm.invoke([HumanMessage(content=prompt)])
            json_match = re.search(r'\{.*\}', response.content, re.DOTALL)
            
            if json_match:
                insights = json.loads(json_match.group())
            else:
                insights = {
                    "main_insights": ["Insights generated"],
                    "trends": ["Trends identified"],
                    "implications": ["Implications noted"],
                    "recommendations": ["Recommendations provided"],
                    "future_research": ["Further study needed"]
                }
            
            state['insights'] = insights
            state['messages'].append(f"✓ Generated {len(insights.get('main_insights', []))} insights")
            
            # Add to timeline
            state['timeline'].append({
                'agent': self.name,
                'action': 'Insights generated',
                'duration': round(time.time() - start_time, 2),
                'timestamp': datetime.now().isoformat()
            })
            
        except Exception as e:
            state['messages'].append(f"⚠ Insight error: {str(e)}")
            state['insights'] = {
                "main_insights": ["Error"],
                "trends": ["Unknown"],
                "implications": ["Unknown"],
                "recommendations": ["Retry"],
                "future_research": ["Complete analysis"]
            }
        
        return state


class ReportBuilderAgent:
    """Agent 5: Builds final report"""
    
    def __init__(self, llm):
        self.llm = llm
        self.name = "Report Builder"
    
    def invoke(self, state: AgentState) -> AgentState:
        """Build report"""
        start_time = time.time()
        
        prompt = f"""You are a Report Builder Agent.

Query: {state['query']}
Findings: {state['analysis'].get('key_findings', [])}
Insights: {state['insights'].get('main_insights', [])}

Write a professional executive summary (4-6 paragraphs):
1. Research question and methodology
2. Key findings with evidence
3. Critical insights and trends
4. Actionable recommendations
5. Implications and future directions

Use academic prose, be specific and evidence-based."""

        try:
            response = self.llm.invoke([HumanMessage(content=prompt)])
            state['final_report'] = response.content
            state['messages'].append("✓ Final report compiled")
            
            # Add to timeline
            state['timeline'].append({
                'agent': self.name,
                'action': 'Report compiled',
                'duration': round(time.time() - start_time, 2),
                'timestamp': datetime.now().isoformat()
            })
            
        except Exception as e:
            state['messages'].append(f"⚠ Report error: {str(e)}")
            state['final_report'] = f"Error: {str(e)}"
        
        return state


# ============= LANGGRAPH WORKFLOW =============

@st.cache_resource
def create_research_graph():
    """Create agent workflow"""
    
    planner = ResearchPlannerAgent(llm)
    rag_searcher = RAGSearchAgent()
    web_searcher = WebSearchAgent(tavily_client)
    analyzer = AnalysisAgent(llm)
    insight_gen = InsightGenerationAgent(llm)
    reporter = ReportBuilderAgent(llm)
    
    workflow = StateGraph(AgentState)
    
    workflow.add_node("plan", planner.invoke)
    workflow.add_node("rag_search", rag_searcher.invoke)
    workflow.add_node("web_search", web_searcher.invoke)
    workflow.add_node("analyze", analyzer.invoke)
    workflow.add_node("insights", insight_gen.invoke)
    workflow.add_node("report", reporter.invoke)
    
    workflow.set_entry_point("plan")
    workflow.add_edge("plan", "rag_search")
    workflow.add_edge("rag_search", "web_search")
    workflow.add_edge("web_search", "analyze")
    workflow.add_edge("analyze", "insights")
    workflow.add_edge("insights", "report")
    workflow.add_edge("report", END)
    
    return workflow.compile()


# ============= AGENT VISUALIZATION =============

def create_agent_visualization():
    """Create visual workflow diagram"""
    
    workflow_html = """
    <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 20px; border-radius: 10px; color: white;">
        <h3 style="text-align: center; margin-bottom: 20px;">Agent Workflow Pipeline</h3>
        <div style="display: flex; justify-content: space-around; align-items: center; flex-wrap: wrap;">
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">1</div>
                <p style="margin-top: 10px;">Planner</p>
            </div>
            <div style="font-size: 30px; margin: 10px;">→</div>
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">2</div>
                <p style="margin-top: 10px;">RAG Search</p>
            </div>
            <div style="font-size: 30px; margin: 10px;">→</div>
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">3</div>
                <p style="margin-top: 10px;">Web Search</p>
            </div>
            <div style="font-size: 30px; margin: 10px;">→</div>
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">4</div>
                <p style="margin-top: 10px;">Analyzer</p>
            </div>
            <div style="font-size: 30px; margin: 10px;">→</div>
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">5</div>
                <p style="margin-top: 10px;">Insights</p>
            </div>
            <div style="font-size: 30px; margin: 10px;">→</div>
            <div style="text-align: center; margin: 10px;">
                <div style="background: white; color: #667eea; padding: 15px; border-radius: 50%; width: 80px; height: 80px; display: flex; align-items: center; justify-content: center; font-weight: bold; margin: 0 auto;">6</div>
                <p style="margin-top: 10px;">Reporter</p>
            </div>
        </div>
    </div>
    """
    return workflow_html


# ============= STREAMLIT UI =============

st.title("🧠 MAAD Researcher")

st.markdown("*Multi-Agent AI Deep(MAAD): RAG • PDF Upload • Citations • Agent Visualization*")

st.markdown("---")

# Sidebar
with st.sidebar:
    st.header("⚙️ System Info")
    st.markdown("""
    <style>
    .tiny-metric .stMetric > label { font-size: 1px !important; }
    .tiny-metric .stMetric > div > div { font-size: 1px !important; }
    </style>
    """, unsafe_allow_html=True)

    # 2-row ultra compact
    col1, col2, col3 = st.columns(3)
    col1.metric("🔍 Searches", len(st.session_state.search_history))
    col2.metric("📄 Docs", len(st.session_state.uploaded_docs))
    col3.metric("🎯 Citations", len(st.session_state.citations))

    col1, col2 = st.columns(2)
    col1.metric("📦 Chunks", f"{db_stats['total_chunks']:,}")
    col2.metric("💾 Size", f"{db_stats['total_size_mb']} MB")
    # col3.metric("📏 Avg Chunk", f"{db_stats['avg_chunk_length']} chars")
    # col4.metric("🔢 Dimensions", db_stats['dimensions'])

    st.info(api_status)
    
    st.markdown("---")
    st.subheader("📤 Upload Research Documents")
    
    uploaded_files = st.file_uploader(
        "Upload PDFs for RAG",
        type=['pdf'],
        accept_multiple_files=True,
        help="Upload research papers, reports, or documents to enhance research"
    )
    
    if uploaded_files:
        if st.button("Process PDFs", type="primary"):
            with st.spinner("Processing PDFs..."):
                all_docs = []
                for uploaded_file in uploaded_files:
                    docs, preview = process_pdf(uploaded_file)
                    if docs:
                        all_docs.extend(docs)
                        st.session_state.uploaded_docs.append({
                            'name': uploaded_file.name,
                            'preview': preview,
                            'chunks': len(docs)
                        })
                
                if all_docs:
                    add_to_vectorstore(all_docs)
                    st.success(f"✅ Processed {len(uploaded_files)} PDFs ({len(all_docs)} chunks)")
    
    # Show uploaded docs
    if st.session_state.uploaded_docs:
        st.markdown("---")
        st.subheader("📚 Uploaded Documents")
        for doc in st.session_state.uploaded_docs:
            with st.expander(doc['name']):
                st.caption(f"Chunks: {doc['chunks']}")
                st.text(doc['preview'][:200] + "...")
    
    st.markdown("---")
    st.markdown("**Agent Pipeline:**")
    agents = [
        "1. 🎯 Research Planner",
        "2. 📄 RAG Search",
        "3. 🔍 Web Search",
        "4. 📊 Critical Analysis",
        "5. 💡 Insight Generator",
        "6. 📝 Report Builder"
    ]
    for agent in agents:
        st.markdown(agent)
    
    st.markdown("---")
    st.markdown("**Cost:** $0.00 🎉")

    # SEARCH HISTORY SIDEBAR
    st.markdown("---")
    st.subheader("📚 Search History")

    if st.session_state.search_history:
        # History list
        for i, entry in enumerate(reversed(st.session_state.search_history[-10:])):  # Last 10
            with st.expander(f"🔍 {entry['query'][:60]}... ({entry['timestamp']})", expanded=False):
                col1, col2, col3 = st.columns([3, 1, 1])
                with col1:
                    st.caption(f"Scope: {entry['scope'].title()}")
                    if entry['domains']:
                        st.caption(f"Domains: {', '.join(entry['domains'])}")
                with col2:
                    if st.button("📋", key=f"load_{entry['id']}", use_container_width=True):
                        st.session_state.results = entry['results']
                        st.session_state.agent_logs = entry['results'].get('messages', [])
                        st.session_state.citations = entry['results'].get('citations', [])
                        st.session_state.agent_timeline = entry['results'].get('timeline', [])
                        st.session_state.current_query = entry['query']
                        st.success(f"Loaded: {entry['query']}")
                        st.rerun()
                with col3:
                    if st.button("🗑️", key=f"del_{entry['id']}", use_container_width=True):
                        st.session_state.search_history = [h for h in st.session_state.search_history if h['id'] != entry['id']]
                        st.rerun()
        
        # Clear all history
        if st.button("🗑️ Clear All History", type="secondary"):
            st.session_state.search_history = []
            st.rerun()
    else:
        st.info("No searches yet. Run a research query to start building history!")


# Main interface
st.markdown("### Enter Your Research Query")

query = st.text_input(
    "Research query",
    value=st.session_state.current_query,
    placeholder="e.g., Impact of AI on healthcare diagnostics",
    key="query_input",
    label_visibility="collapsed"
)

col1, col2 = st.columns([2, 1])
with col1:
    research_button = st.button("🔬 Start Research", type="primary", use_container_width=True)
with col2:
    if st.session_state.results:
        if st.button("🔄 Clear", use_container_width=True):
            st.session_state.results = None
            st.session_state.agent_logs = []
            st.session_state.citations = []
            st.session_state.agent_timeline = []
            st.rerun()

# Agent Visualization
if not st.session_state.results:
    st.markdown("---")
    st.markdown("### 🎯 Agent Workflow Visualization")
    st.markdown(create_agent_visualization(), unsafe_allow_html=True)

# Execute research
if research_button and query:
    st.session_state.agent_logs = []
    st.session_state.results = None
    st.session_state.citations = []
    st.session_state.agent_timeline = []
    
    with st.spinner("🤖 AI Agents researching..."):
        try:
            initial_state = {
                "query": query,
                "research_plan": {},
                "search_results": [],
                "rag_results": [],
                "analysis": {},
                "insights": {},
                "final_report": "",
                "messages": [],
                "citations": [],
                "timeline": []
            }
            
            st.info("🎯 Executing multi-agent pipeline...")
            
            graph = create_research_graph()
            final_state = graph.invoke(initial_state)
            
            st.session_state.results = final_state
            st.session_state.agent_logs = final_state['messages']
            st.session_state.citations = final_state.get('citations', [])
            st.session_state.agent_timeline = final_state.get('timeline', [])
            
            st.success("✅ Research complete!")
            st.balloons()
            

            # SAVE TO HISTORY
            history_entry = {
                'id': len(st.session_state.search_history),
                'query': query,
                'scope': st.session_state.search_scope,
                'domains': st.session_state.web_domains.copy(),
                'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M'),
                'results': final_state.copy(),  # Full results
                'summary_preview': final_state.get('final_report', '')[:200] + '...'
            }
            st.session_state.search_history.append(history_entry)
            st.rerun()

            
        except Exception as e:
            st.error(f"❌ Error: {str(e)}")

# Display results
if st.session_state.results:
    results = st.session_state.results
    
    # Agent Timeline
    if st.session_state.agent_timeline:
        with st.expander("⏱️ Agent Execution Timeline", expanded=False):
            for event in st.session_state.agent_timeline:
                col1, col2, col3 = st.columns([2, 2, 1])
                with col1:
                    st.markdown(f"**{event['agent']}**")
                with col2:
                    st.markdown(f"{event['action']}")
                with col3:
                    st.markdown(f"`{event['duration']}s`")
            
            total_time = sum(e['duration'] for e in st.session_state.agent_timeline)
            st.info(f"**Total Execution Time:** {total_time:.2f}s")
    
    # Agent Logs
    with st.expander("🤖 Agent Activity Logs", expanded=False):
        for log in st.session_state.agent_logs:
            st.text(log)
    
    st.markdown("---")
    
    # Metrics
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.metric("Web Sources", len(results.get('search_results', [])))
    with col2:
        st.metric("Chunks Retrieved", len(results.get('rag_results', [])))
    with col3:
        confidence = results.get('analysis', {}).get('confidence_level', 'medium')
        st.metric("Confidence", confidence.title())
    with col4:
        st.metric("Citations", len(st.session_state.citations))
    
    st.markdown("---")
    
    # Research Plan
    st.subheader("🎯 Research Strategy")
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("**Topics:**")
        for topic in results.get('research_plan', {}).get('topics', []):
            st.markdown(f"• {topic}")
    
    with col2:
        st.markdown("**Queries:**")
        for q in results.get('research_plan', {}).get('search_queries', [])[:5]:
            st.markdown(f"• {q}")
    
    st.markdown("---")
    
    # RAG Results
    if results.get('rag_results'):
        st.subheader("📄 Uploaded Document Insights")
        for i, rag in enumerate(results.get('rag_results', [])[:3], 1):
            with st.expander(f"Document {i}: {rag.get('source', 'Unknown')} (Page {rag.get('page', 'N/A')})"):
                st.text(rag.get('content', '')[:500] + "...")
        st.markdown("---")
    
    # Web Sources
    st.subheader("📚 Web Research Sources")
    for i, source in enumerate(results.get('search_results', [])[:10], 1):
        with st.expander(f"Source {i}: {source.get('title', 'Untitled')}"):
            st.markdown(f"**URL:** [{source.get('url', 'N/A')}]({source.get('url', '#')})")
            st.markdown(f"**Relevance:** {source.get('score', 0):.2%}")
            st.text(source.get('content', '')[:500] + "...")
    
    st.markdown("---")
    
    # Citations
    if st.session_state.citations:
        st.subheader("📑 Citations & Sources")
        
        # Group by agent
        citations_by_agent = {}
        for citation in st.session_state.citations:
            agent = citation['agent']
            if agent not in citations_by_agent:
                citations_by_agent[agent] = []
            citations_by_agent[agent].append(citation)
        
        for agent, citations in citations_by_agent.items():
            with st.expander(f"{agent} - {len(citations)} citations"):
                for i, cite in enumerate(citations[:10], 1):
                    st.markdown(f"**{i}.** [{cite['source_title']}]({cite['source_url']})")
                    st.caption(cite['claim'])
        
        st.markdown("---")
    
    # Analysis
    st.subheader("🔬 Critical Analysis")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("**Key Findings:**")
        for finding in results.get('analysis', {}).get('key_findings', []):
            st.markdown(f"• {finding}")
    
    with col2:
        st.info(f"**Quality:**\n{results.get('analysis', {}).get('data_quality', 'N/A')}")
        
        gaps = results.get('analysis', {}).get('gaps', 'None')
        if gaps and isinstance(gaps, str) and gaps.lower() != 'none':
            st.warning(f"Gaps: {gaps}")
        elif isinstance(gaps, list) and gaps:
            st.warning(f"Gaps: {', '.join(gaps[:3])}...")
        # if gaps.lower() != 'none':
        #     st.warning(f"**Gaps:**\n{gaps}")
    
    st.markdown("---")
       
    # Insights
    st.subheader("💡 Insights & Recommendations")
    
    tab1, tab2, tab3, tab4 = st.tabs(["Insights", "Trends", "Recommendations", "Future Research"])
    
    with tab1:
        for insight in results.get('insights', {}).get('main_insights', []):
            st.markdown(f"✓ {insight}")
    
    with tab2:
        for trend in results.get('insights', {}).get('trends', []):
            st.markdown(f"📈 {trend}")
    
    with tab3:
        for rec in results.get('insights', {}).get('recommendations', []):
            st.markdown(f"→ {rec}")
    
    with tab4:
        for area in results.get('insights', {}).get('future_research', []):
            st.markdown(f"🔍 {area}")
    
    st.markdown("---")
    
    # # Executive Summary
    # COLLAPSIBLE EXECUTIVE SUMMARY
    with st.expander("📝 Executive Summary", expanded=True):
        st.markdown(results.get('final_report', ''))

    st.markdown("---")

    # ===== PLAGIARISM CHECKER =====
    st.markdown("---")
    st.subheader("🔍 Plagiarism Checker")

    col_btn1, col_btn2 = st.columns([3, 1])

    with col_btn1:
        if st.button("🚀 Check Plagiarism", type="secondary", use_container_width=True):
            st.session_state.plagiarism_active = True

    with col_btn2:
        st.caption("Powered by Winston AI")

    if st.session_state.get('plagiarism_active', False):
        st.markdown("---")
        
        # Text input
        plagiarism_text = st.text_area(
            "Paste content to check:",
            height=200,
            placeholder="Paste your text here... (Max 500 words free)",
            help="Winston AI checks AI-generated content vs human writing"
        )
        
        word_count = len(plagiarism_text.split())
        
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Words", f"{word_count:,}")
        
        if st.button("✅ Analyze Now", type="primary"):
            if word_count == 0:
                st.warning("Please paste some text first!")
            elif word_count > 500:
                st.error("""
    ## 🚫 Premium Required

    **Detected**: {word_count:,} words

    **Free limit**: 500 words  
    **Upgrade**: [Winston AI Premium](https://gowinston.ai/pricing)

    **Why?** Large documents require paid API access.
                """.format(word_count=word_count))
            else:
                with st.spinner("Analyzing with Winston AI..."):
                    result = check_plagiarism(plagiarism_text, winston_api_key)
                    # st.json(result)
                    analysis = analyze_plagiarism_result(result, word_count)
                    st.markdown(analysis)
                    
                    # Results storage
                    st.session_state.plagiarism_result = {
                        'text': plagiarism_text,
                        'word_count': word_count,
                        'result': result,
                        'timestamp': datetime.now().strftime('%H:%M')
                    }
        
        # Show last result
        if st.session_state.get('plagiarism_result'):
            with st.expander("📋 Last Check", expanded=False):
                last = st.session_state.plagiarism_result
                st.info(f"**{last['word_count']:,} words** • {last['timestamp']}")


    # HISTORY TAB
    tab1, tab2 = st.tabs(["📊 Current Results", "📚 Past Searches"])

    with tab1:
        # Your existing results sections (Critical Analysis, Insights, Downloads, etc.)
        pass  # Keep all your current result displays here

    with tab2:
        st.subheader("📚 Previous Searches")
        if st.session_state.search_history:
            for entry in reversed(st.session_state.search_history[-5:]):  # Show last 5
                with st.expander(f"🔍 {entry['query']} ({entry['timestamp']})"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.markdown(f"**Scope**: {entry['scope'].title()}")
                        if entry['domains']:
                            st.markdown(f"**Domains**: {', '.join(entry['domains'])}")
                        st.caption(f"**Summary**: {entry['summary_preview']}")
                    with col2:
                        st.markdown("**Quick Actions**")
                        col_btn1, col_btn2 = st.columns(2)
                        with col_btn1:
                            if st.button("📋 Load Full Results", key=f"load_tab_{entry['id']}"):
                                st.session_state.results = entry['results']
                                st.session_state.current_query = entry['query']
                                st.rerun()
                        with col_btn2:
                            if st.button("📥 Export Summary", key=f"exp_{entry['id']}"):
                                summary = entry['results'].get('final_report', '')
                                st.download_button(
                                    "PDF", export_to_pdf(summary, entry['query']),
                                    f"{entry['query'][:30]}.pdf", "application/pdf"
                                )
        else:
            st.info("No previous searches. Run queries to build history!")


    

    # ===== UPDATED EXPORT SECTION - Short=300 words + Full =====
    st.markdown("---")
    st.subheader("📥 Downloading Options")

    col1, col2, col3 = st.columns([1, 1, 1])

    with col1:
        st.markdown("**📄 Summary Length**")
        summary_type = st.radio("Choose:", ["Short (300 words)", "Full"], horizontal=True)
        
        if summary_type == "Short (300 words)":
            
            # 1. FIRST: Remove title/headers
            full_report = results.get('final_report', '')
            lines = full_report.split('\n')
            clean_lines = []
            
            # Skip title/headers (first 2-3 lines usually)
            for i, line in enumerate(lines[:8]):  # Check more lines
                line = line.strip()
                if line and not line.startswith('**') and not line.startswith('Executive') and len(line) > 20:
                    clean_lines.append(line)
                    break  # Take only first real content line
            
            # 2. Build first full paragraph (complete, no cuts)
            cleaned_report = '\n'.join(clean_lines)
            
            # Find COMPLETE first paragraph (until double newline OR natural break)
            if '\n\n' in cleaned_report:
                first_para = cleaned_report.split('\n\n')[0].strip()
            else:
                # No double newline - take until first long pause (2+ empty lines or 800 chars max)
                first_para = cleaned_report.split('\n\n', 1)[0].strip()
                if len(first_para) > 800:  # Max safety limit
                    # Find natural sentence break
                    sentences = first_para.split('. ')
                    first_para = '. '.join(sentences[:3]) + '.'  # First 3 sentences
            
            # Clean markdown only
            first_para = first_para.replace('**', '').replace('*', '').strip()
            
            export_content = f"""EXECUTIVE SUMMARY: {query.upper()}

        {first_para}

        ---
        MAAD Researcher | {datetime.now().strftime('%Y-%m-%d %H:%M')}"""
            
            export_title = f"Exec-Summary-{query[:25]}"



            
        else:  # Full
            export_content = results.get('final_report', '')
            export_title = f"Full-Report-{query[:25]}"
        
        st.text_area("Preview:", export_content, height=150, label_visibility="collapsed")

    with col2:
        st.markdown("**💾 Export Format**")
        exp_format = st.radio("Format:", ["PDF", "PPT", "TXT"], horizontal=True)
        
        if exp_format == "PDF":
            pdf_data = export_to_pdf(export_content, export_title)
            st.download_button(
                "📄 PDF", pdf_data, f"{export_title}.pdf", 
                "application/pdf", use_container_width=True
            )
        elif exp_format == "PPT":
            ppt_data = export_to_pptx(export_content, export_title)
            st.download_button(
                "📊 PPT", ppt_data, f"{export_title}.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )
        else:
            st.download_button(
                "📝 TXT", export_content, f"{export_title}.txt",
                "text/plain", use_container_width=True
            )

    with col3:
        st.markdown("**💼 Share**")
        
        # LinkedIn post
        linkedin_post = f"""🚀 {query.upper()}\n\n{export_content}...\n\n#AI #Research #Tech #DataScience"""
        if st.button("💼 Copy LinkedIn Post", use_container_width=True):
            st.code(linkedin_post, language="markdown")
            st.success("✅ Copy the above & paste to LinkedIn!")

    # Keep existing JSON/TXT downloads
    # col1, col2 = st.columns(2)
    # jsondata = json.dumps(results, indent=2, default=str)
    # col1.download_button("JSON Report", jsondata, f"report-{query[:20]}.json", "application/json")
    # textreport = f"RESEARCH REPORT\nQuery: {query}\n\nEXECUTIVE SUMMARY\n{results.get('finalreport', '')}"
    # col2.download_button("Text Summary", textreport, f"summary-{query[:20]}.txt", "text/plain")


#     # Download
#     col1, col2, col3 = st.columns(3)
    
#     with col1:
#         json_data = json.dumps(results, indent=2, default=str)
#         st.download_button(
#             "📥 JSON Report",
#             json_data,
#             f"report_{query[:20]}.json",
#             "application/json"
#         )
    
#     with col2:
#         text_report = f"""RESEARCH REPORT
# ================
# Query: {query}

# EXECUTIVE SUMMARY:
# {results.get('final_report', '')}

# KEY FINDINGS:
# {chr(10).join(['• ' + f for f in results.get('analysis', {}).get('key_findings', [])])}

# RECOMMENDATIONS:
# {chr(10).join(['• ' + r for r in results.get('insights', {}).get('recommendations', [])])}

# CITATIONS: {len(st.session_state.citations)}
# Sources: {len(results.get('search_results', []))} web + {len(results.get('rag_results', []))} RAG
# """
#         st.download_button(
#             "📄 Text Summary",
#             text_report,
#             f"summary_{query[:20]}.txt",
#             "text/plain"
#         )
    
#     with col3:
#         # Citations export
#         citations_text = "CITATIONS\n==========\n\n"
#         for i, cite in enumerate(st.session_state.citations, 1):
#             citations_text += f"{i}. {cite['source_title']}\n"
#             citations_text += f"   URL: {cite['source_url']}\n"
#             citations_text += f"   Agent: {cite['agent']}\n\n"
        
#         st.download_button(
#             "📑 Citations",
#             citations_text,
#             f"citations_{query[:20]}.txt",
#             "text/plain"
#         )

# Example queries
if not st.session_state.results:
    st.markdown("---")
    st.subheader("💡 Example Queries")
    
    examples = [
        "Analyze quantum computing breakthroughs and cybersecurity implications",
        "Impact of AI on healthcare diagnostics with ethical considerations",
        "Current state of renewable energy storage technology",
        "Gene editing advances and CRISPR ethical concerns",
        "Remote work effects on productivity and mental health",
        "Nuclear fusion energy recent developments"
    ]
    
    cols = st.columns(2)
    for idx, example in enumerate(examples):
        with cols[idx % 2]:
            if st.button(example, key=f"ex_{idx}", use_container_width=True):
                st.session_state.current_query = example
                st.rerun()

st.markdown("---")
st.caption("🚀 Advanced Multi-Agent Research System | Groq + Tavily + RAG + ChromaDB")